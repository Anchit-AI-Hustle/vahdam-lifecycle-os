#!/usr/bin/env python3
"""Build the VAHDAM Lifecycle OS pitch deck as brand-styled PPTX + print HTML
(the HTML is rendered to PDF by Chromium). One slide-data source -> both."""
import html
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG   = RGBColor(0x06,0x13,0x0D)   # near-black forest canvas
PANEL= RGBColor(0x0F,0x1D,0x18)
GOLD = RGBColor(0xAB,0x87,0x43)
CREAM= RGBColor(0xFB,0xF5,0xEA)
MUT  = RGBColor(0xCD,0xD8,0xD2)

# ---- Slide data (title, eyebrow, and either bullets or a table) ----
S = [
 {"eyebrow":"VAHDAM India · Internal Platform → SaaS Candidate","title":"Lifecycle OS — Smart Brain",
  "lead":"One system that reads our data, watches every competitor, plans the next 15 days, generates every asset (mailers, Google/Meta/TikTok ads, landing pages) and talks to customers in a human voice. Humans verify; the brain executes.",
  "bullets":["Live at /brain · DB-linked · daily cron · phase-2 live-push excluded by design",
             "Built on the existing shell · Supabase · 6-provider LLM cascade · ~$0 infra",
             "Now includes the V2 Lifecycle OS layer (2026-07): Mailer Calendar, UK Non-Engagers, Social OS"]},
 {"eyebrow":"Why · bandwidth & resource utilisation","title":"The marketing team's week, before and after",
  "table":[["Task","Manual today","With Smart Brain","Saved"],
           ["Performance digging","~6 h/wk in sheets","Daily auto-analysis + chat console","~5.5 h"],
           ["Calendar planning","~4 h/wk, gut-driven","15-day rolling, self-reviewing daily","~3.5 h"],
           ["Mailer + ad + LP creation","~2 days per set","Full funnel in <60s, on-brand by construction","~85%"],
           ["Competitor monitoring","Ad-hoc inbox digging","Continuous isolated stream + signals","~2 h"],
           ["Customer education","Headcount-limited","Voice agents 24x7, every market","scale"]]},
 {"eyebrow":"What it is · module map","title":"Six modules, clean contracts",
  "bullets":["1 · Knowledge Base — creative-level index of what worked",
             "2 · Analysis Engine — thresholds, cohorts, daily sync",
             "3 · Competitor Stream — isolated; advisory signals only (never qualifies our winners)",
             "4 · Calendar Intelligence — 15-day rolling, festival-aware, self-reviewing",
             "5 · Generation Engine — mailers, Google/Meta/TikTok ads, landing pages as platform-ready objects",
             "6 · Human Review — every campaign verified at launch; hard weekly recalibration gate"]},
 {"eyebrow":"Flows · the operating loop","title":"Daily is automated. Weekly is human.",
  "bullets":["Daily 03:30 UTC cron, no humans: re-extract festivals & peaks from history",
             "Full re-analysis vs the DB — never follow the calendar blindly, update it",
             "Apply human feedback, swap weak angles, extend the 15-day horizon",
             "Fold MVT winners into learned weights + competitor signals",
             "Generate complete asset sets for approved slots due within 3 days",
             "HUMAN GATE: approve → final; reject → reasons feed back. Recalibration >7 days old hard-blocks auto-approval"]},
 {"eyebrow":"Feature · the LLM-style console (ChaiGPT)","title":"Operate it like ChatGPT, grounded in live data",
  "bullets":["Claude-style chat that plans multi-step work and calls real tools",
             "Collapsible tool trace: every call, its args, and results",
             "Evidence contract: quote exact tool-sourced figures, name the metric + expected impact, state a hypothesis, quote competitor benchmarks",
             "Works across the entire 6-provider waterfall including free tiers"]},
 {"eyebrow":"Feature · calendar intelligence","title":"A 15-day plan that argues for itself",
  "bullets":["Each slot: date × market × cohort (live size) × channel × theme/angle × hero product",
             "Every entry carries a rationale and a confidence score",
             "Festival-weighted and send-time tuned; MVT loop records per-variant outcomes",
             "Daily review narrates what changed and why (angle swaps, skips from your feedback)"]},
 {"eyebrow":"Feature · generation engine","title":"One approval → a complete funnel",
  "bullets":["Approved slot (cohort + angle + theme) fans out to creative, landing, follow-up",
             "HTML mailer (brand-locked: palette · Lao MN · banned-phrase scrub)",
             "Google RSA · Meta · TikTok copy + baked-text creatives",
             "Two landing pages (for the mailer and for the ads) hosted at /lp/:id",
             "Audience specs (retargeting 14d, lookalike 1-3%) — stored as platform-ready JSON; phase-2 push needs zero refactor"]},
 {"eyebrow":"Feature · voice agents","title":"Telecalling held ROAS 3.5-4. This reproduces it at ~$0.",
  "bullets":["Value justification is conversational: duration, benefits, per-cup math",
             "Brand / product / collection agents — voice in, voice out, honest by design",
             "One script tag on Shopify routes to the right specialist; mailers & LPs deep-link to /agent",
             "Knowledge = catalog slice + scraped official product pages + brand KB"]},
 {"eyebrow":"Trust · accuracy criteria & improvement","title":"How we know it's right — and how it improves",
  "table":[["Surface","Acceptance criteria (launch)","Improvement loop"],
           ["Analysis","Scores reproducible from raw metrics; pass-rate 20-60%","Weekly recalibration adjusts thresholds"],
           ["Calendar","≥80% of tentative slots approved without edits by wk 2","Rejections suppress angles; MVT winners boost weights"],
           ["Assets","100% brand-compliant by code; ≥70% approved w/o edits","Rejection notes fed into next prompt"],
           ["Voice","0 medical claims · 0 banned phrases · ≥85% resolution","Persona + knowledge updates from transcript audit"],
           ["System","Daily cron ≥99%; recalibration never >7d","Confidence relaxes review load only on proof"]]},
 {"eyebrow":"Impact hypothesis","title":"What we expect it to move",
  "bullets":["Bandwidth: about −70% production hours per campaign set (team shifts to verification + strategy)",
             "Send-to-plan latency: 15 days → daily, re-validated every morning",
             "Creative win-rate: +20-30% (only threshold-passing angles seed new campaigns; MVT compounds)",
             "Education channel ROAS: ~3.5x via voice agents",
             "Cost: ~$0 (free-tier DB + existing Vercel + free LLM tiers); Risk: gated (human-verified, weekly recalibration)"]},
 {"eyebrow":"Go-live plan","title":"From today to autopilot",
  "bullets":["Week 0 (done): Smart Brain live at /brain; cron + review queue + agents ready",
             "Week 1: point the linked DB at real Shopify/Klaviyo/ads exports; harden RLS; first recalibration",
             "Week 2-3: shadow mode — brain plans & generates, team ships manually; tune thresholds weekly",
             "Week 4+: phase 2 — Klaviyo + ad-platform push on stored objects; performance flows back; loop closed"]},
 {"eyebrow":"Where this goes","title":"VAHDAM first. SaaS later.",
  "bullets":["Nothing in the brain is VAHDAM-specific by code — brand kit, thresholds, catalog, agents are DB rows",
             "Multi-tenancy is one linked DB per brand away",
             "Wedge: D2C brands with scattered history, no analyst bandwidth, a price-justification problem",
             "Pricing instinct: per-brand platform fee + per-generated-campaign / per-agent-conversation usage"]},
 {"eyebrow":"What's new · V1 → V2 (2026-07)","title":"The current-state layer",
  "bullets":["V1/V2 taxonomy adopted (2026-07-03): V1 = legacy base; V2 = Lifecycle OS additions",
             "V2: cohort Mailer Calendar, UK Non-Engagers hub, tier-routed cascades + video-core, Social Media OS, retention/influencer KB",
             "Access: the demo/mock gate was removed — every signed-in user gets full live access",
             "Design: UI locked to a single forest-green theme (dark/dusk/light switcher removed)",
             "Ops: domain + OAuth migration tooling (migrate-domains / migrate-oauth) keeps Google sign-in working on new domains"]},
 {"eyebrow":"Live","title":"See it running",
  "bullets":["Console: /brain     ·     Operator: /chaigpt     ·     Buyer agent: /agent",
             "This deck: /deck     ·     PRD: /prd-deck",
             "Everything on the brand's forest-green canvas, gold + cream type, Lao MN / Proxima Nova"]},
]

# ---------------- PPTX ----------------
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
def bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color
def box(slide, l,t,w,h):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; return tb,tf

for s in S:
    sl = prs.slides.add_slide(blank); bg(sl, BG)
    _,tf = box(sl,0.9,0.55,11.5,0.5)
    r = tf.paragraphs[0].add_run(); r.text = s["eyebrow"].upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=GOLD; r.font.name="Arial"
    _,tf = box(sl,0.9,1.05,11.5,1.2)
    r = tf.paragraphs[0].add_run(); r.text=s["title"]
    r.font.size=Pt(34); r.font.bold=True; r.font.color.rgb=CREAM; r.font.name="Georgia"
    y=2.4
    if s.get("lead"):
        _,tf=box(sl,0.9,y,11.5,1.1); p=tf.paragraphs[0]; rr=p.add_run(); rr.text=s["lead"]
        rr.font.size=Pt(14); rr.font.color.rgb=MUT; rr.font.name="Arial"; y+=1.2
    if s.get("bullets"):
        _,tf=box(sl,0.9,y,11.5,7.4-y-0.4)
        for k,b in enumerate(s["bullets"]):
            p = tf.paragraphs[0] if k==0 else tf.add_paragraph()
            rr=p.add_run(); rr.text="•  "+b; rr.font.size=Pt(15); rr.font.color.rgb=CREAM; rr.font.name="Arial"
            p.space_after=Pt(8)
    if s.get("table"):
        data=s["table"]; rows=len(data); cols=len(data[0])
        gt = sl.shapes.add_table(rows,cols,Inches(0.9),Inches(y),Inches(11.5),Inches(0.5*rows)).table
        for ci in range(cols):
            for ri in range(rows):
                cell=gt.cell(ri,ci); cell.text=data[ri][ci]
                para=cell.text_frame.paragraphs[0]; para.runs[0].font.size=Pt(10.5); para.runs[0].font.name="Arial"
                if ri==0:
                    cell.fill.solid(); cell.fill.fore_color.rgb=FOREST if False else RGBColor(0x00,0x4A,0x2B)
                    para.runs[0].font.color.rgb=CREAM; para.runs[0].font.bold=True
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb=PANEL; para.runs[0].font.color.rgb=MUT
    # page number
    _,tf=box(sl,12.2,7.0,1.0,0.4); r=tf.paragraphs[0].add_run(); r.text=str(S.index(s)+1)
    r.font.size=Pt(10); r.font.color.rgb=GOLD
prs.save("docs/deck.pptx"); print("wrote docs/deck.pptx")

# ---------------- print HTML (-> chromium PDF) ----------------
def esc(t): return html.escape(t)
slides_html=[]
for idx,s in enumerate(S):
    inner = f'<div class="eyebrow">{esc(s["eyebrow"])}</div><h1>{esc(s["title"])}</h1>'
    if s.get("lead"): inner += f'<p class="lead">{esc(s["lead"])}</p>'
    if s.get("bullets"): inner += "<ul>" + "".join(f"<li>{esc(b)}</li>" for b in s["bullets"]) + "</ul>"
    if s.get("table"):
        d=s["table"]; inner += "<table><thead><tr>" + "".join(f"<th>{esc(c)}</th>" for c in d[0]) + "</tr></thead><tbody>"
        for r in d[1:]: inner += "<tr>" + "".join(f"<td>{esc(c)}</td>" for c in r) + "</tr>"
        inner += "</tbody></table>"
    inner += f'<div class="pg">{idx+1} / {len(S)}</div>'
    slides_html.append(f'<section class="slide">{inner}</section>')
page = """<!doctype html><html><head><meta charset="utf-8"><style>
@page { size: 13.333in 7.5in; margin:0; }
* { box-sizing:border-box; }
body { margin:0; font-family:'Helvetica Neue',Arial,sans-serif; }
.slide { width:13.333in; height:7.5in; padding:0.55in 0.9in; background:#06130d;
  background-image: radial-gradient(1100px 600px at 80% -10%, rgba(171,135,67,.12), transparent);
  color:#FBF5EA; page-break-after:always; position:relative; overflow:hidden; }
.eyebrow { font-size:12pt; letter-spacing:.18em; text-transform:uppercase; color:#AB8743; font-weight:700; }
h1 { font-family:Georgia,serif; font-size:30pt; margin:10pt 0 12pt; color:#FBF5EA; line-height:1.12; }
.lead { font-size:14pt; color:#cdd8d2; max-width:10.5in; line-height:1.6; }
ul { margin:6pt 0 0; padding-left:20pt; }
li { font-size:14.5pt; color:#FBF5EA; line-height:1.5; margin:7pt 0; max-width:11in; }
table { border-collapse:collapse; width:100%; margin-top:8pt; font-size:11pt; }
th { background:#004A2B; color:#FBF5EA; text-align:left; padding:7pt 9pt; }
td { background:#0f1d18; color:#cdd8d2; border:1px solid rgba(171,135,67,.25); padding:7pt 9pt; }
.pg { position:absolute; bottom:0.35in; right:0.6in; color:#AB8743; font-size:10pt; }
</style></head><body>""" + "".join(slides_html) + "</body></html>"
open("scratchpad/deck_print.html","w",encoding="utf-8").write(page); print("wrote scratchpad/deck_print.html")
